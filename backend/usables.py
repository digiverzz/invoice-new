#by rk
import os
import threading
import creds
import base64
from datetime import datetime
import docx
import pytesseract
import io
from creds import mail
import email
from PIL import Image
from pptx import Presentation
from pdf2image import convert_from_bytes
from predict_data import predict
from datetime import datetime
import json
import pymongo
from urllib.request import urlopen
import hashlib
from difflib import SequenceMatcher
#comment down while deploying
# pytesseract.pytesseract.tesseract_cmd = r'D:\Tesseract-OCR\tesseract.exe'

Files = {}

def IsValidUser(name):
    if creds.collection.find_one({"uid":name}):
        return True
    return False

""" def file_compare(uploadFile,presetFile):
        result = 0
        f1=hashlib.sha1()

        f1.update(uploadFile.read())

        for file in presetFile:
            f2=hashlib.sha1()
            f2.update(file)
            result = max(SequenceMatcher(None,f1.hexdigest(),f2.hexdigest()).ratio(),result)

        return result """

def file_compare(uploadFile,presetFile):
        result = 0.0
        f1=hashlib.sha1()

        f1.update(uploadFile)

        for file in presetFile:
            f2=hashlib.sha1()
            f2.update(file)
            result = max((SequenceMatcher(None,f1.hexdigest(),f2.hexdigest()).ratio())*100,result)
            
        if(result>=20.0):
          print("choosen by function is bestnafaur.pt")
          return "bestnafaur.pt"
        else:
          print("choosen by function is best_1000.pt")
          return "best_1000.pt"


def pathtobytes(dir_path):
    res = []

    for i in os.listdir(dir_path)[1:]:
        byte = open(os.path.join(dir_path,i),"rb")
        res.append(byte.read())
    return res

def getIndex(name):
    if not creds.es.indices.exists(index=name):
        creds.es.indices.create(index=name)
    return name


def get_files(username):
    global Files
    files = Files.get(username,False)
    if files:
        return files
    else:
        Files[username] = creds.es.search(index=username,query={"match_all":{}})['hits']['hits']
        return Files[username]




def get_docx_Text(file_bytes):
    doc = docx.Document(io.BytesIO(file_bytes))
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)

def get_image_Text(file_bytes):
    image = Image.open(io.BytesIO(file_bytes))
    context =  pytesseract.image_to_string(image)
    return context

def get_pdf_Text(file_bytes):
    #images = convert_from_bytes(file_bytes,poppler_path=popplers_path)
    images = convert_from_bytes(file_bytes)
    return "\n".join([pytesseract.image_to_string(img) for img in images])
    
def get_ppt_Text(file_bytes):
    ppt = Presentation(io.BytesIO(file_bytes))
    return "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes])


def extractor(bytes,ext):
    type_functions = {
        "docx": get_docx_Text,
        "pdf": get_pdf_Text,
        "pptx": get_ppt_Text,
        "txt": lambda txt_bytes: txt_bytes.decode()

    }

    func = type_functions.get(ext,get_image_Text)

    return func(bytes)


def data_url_to_image(data_url,filename):

        ext = filename.split('.')[-1]
        raw_b64 = data_url.split(",")[1]

        #base64 to bytes
        file_byte = base64.b64decode(raw_b64)
        
        context = extractor(file_byte,ext)
        return context



def getdatetime():
    now = datetime.now()
    return now



def fetchMail():
    type, data = mail.search(None, 'ALL')
    mail_ids = data[0]
    id_list = mail_ids.split()[::-1]
    mailbox = []

    for msgnum in id_list[::-1]:
        # fetch the message by its ID
        typ, msg_data = mail.fetch(msgnum, '(RFC822)')
        raw_msg = email.message_from_bytes(msg_data[0][1])
        mailObj = {"from":raw_msg["From"],"subject":raw_msg["Subject"],"contents":[]} 

        for part in raw_msg.walk():

            if part.get_content_maintype() == 'multipart':
                continue

            if part.get('Content-Disposition') is None:
                continue
            # extract the filename and content of the attachment
            filename = part.get_filename()
            content = part.get_payload(decode=True)
            # #print(filename)
            
            if filename.endswith('.png') or filename.endswith(".jpeg") or filename.endswith(".jpg") and content is not None:

                mailObj["contents"].append({"filename":filename,"content":content})
            
            elif filename.endswith('.pdf') or filename.endswith(".docx") and content is not None:

                mailObj["contents"].append({"filename":filename,"content":content})

        mailbox.append(mailObj)
    
    #print([i['from'] for i in mailbox])

    return mailbox


    

async def status_update(filename:str,uid:str,status:str):
    
    
    indexname = getIndex(uid)


    query = {
    "match_phrase":{
    "filename":filename
    }
    }
    try:
        id = creds.es.search(index=indexname,query=query)['hits']['hits'][0]['_id']

        result = await creds.es.update(index=indexname,id=id,doc={"status":status})
    except:
        result=[]
        pass    
    return result



async def elastic_upload(username:str,status:str,filename:list,Size:list,dataurl:list):

    # #print(data)
    indexname = getIndex(username)
    res = []
    for fname,size,url in zip(filename,Size,dataurl):
        # #print(fname)

        context = data_url_to_image(url,fname)
        format = {
            "username":username,
            "filename":fname,
            "size":size,
            "context":context,
            "datetime":getdatetime(),
            "dataurl":url,
            "status":status
        }

        res.append(creds.es.index(index=indexname,document=format))
        # #print(res[-1])

    return {"result":res}


def predict_from_mail():
    data = fetchMail()
    response = []
    uri = "mongodb+srv://digiverz:digiverz@cluster0.ngqcelw.mongodb.net/?retryWrites=true&w=majority"
    client = pymongo.MongoClient(uri)
    database = client['invoice']
    #print("data",len(data))
    try:
        for i in range(len(data)):
            emailid = data[i]['from'].split("<")[-1].replace(">","")
            collection = database['users']
            user_data = collection.find_one({"email":emailid},{"_id":0,"id":1,"email":1,"dept":1,"name":1,"uid":1,"role":1,"status":1})    
            #print(user_data,i)
            if user_data!="None" or type(user_data)!=None or user_data!='':
                for j in range(len(data[i]['contents'])):
                    d = {
            "response":"",
                }
                    data_res = predict(data[i]['contents'][j]['content'],"english")
                    d['response'] = json.dumps(data_res)
                    if user_data["role"] == "Associate Practice Lead":
                        user_data["l1"] = "yes"
                        user_data["l2"] = "no"
                    if user_data["role"] == "Employee" :
                        user_data["l1"] = "no"                
                        user_data["l2"] =  "no"

                    if user_data["role"] == "Practice Lead":
                        user_data["l2"] = "yes"
                        user_data["l1"] = "yes"

                    user_data["l3"] = "no"
                    user_data['submitted'] = datetime.today().strftime('%d/%m/%Y')
                    z = {**d, **user_data}
                    response.append(z)
                    #print("__________________________________________________________")
                
        else:
            pass
        
    except TypeError:
        pass

    collection1 = database['request']
    for i in range(len(response)):
        collection1.insert_one(response[i])


def UploadFileHdfs(file,filename):
    path = os.getenv('hdfspath')
    def do():
        creds.hdfs.create_file(path+f'/{filename}',file)
        print(f"{filename} saved in hdfs..")
        event.set()
    
    event = threading.Event()
    thread = threading.Thread(target=do)
    
    thread.start()
    event.wait(timeout=2)
    if not event.is_set():
        print(f"Failed to save {filename} in hdfs..")


def dataurltobytes(url):
    return urlopen(url).read()


####################################### for dashboard and size calculation ###################################
#size calculator
def size_calculator(b64string):
    return (len(b64string) * 3) / 4 - b64string.count('=', -2)

#size converter
def size_converter(value,unit,expected="KB"):
    unit_map = {
        "bytes":{
        "MB":lambda v: v/1000000,
        "GB":lambda v: v/1000000000,
        "KB":lambda v: v/1000
    }}

    return unit_map.get(unit,{}).get(expected,lambda v: v)(value)


#for mongo
def update_size(username,size,type):
      
      myquery = { "uid" : username }
      doc = creds.collection.find_one(myquery, {'_id': 0,"used_size":1,"total_size":1})
      myquery = { "used_size": doc['used_size'] }
      if doc['used_size']<doc['total_size']:
        if type=="1":
            newsize = size+doc['used_size']
            newsizedata = { "$set": { "used_size": newsize } }
            creds.collection.update_one(myquery,newsizedata)
            return {"Newupdatedsize": newsize,"Unit": "MB","status":"success"}
        elif type=="-1":
            newsize = doc['used_size']-size
            newsizedata = { "$set": { "used_size": newsize } }
            creds.collection.update_one(myquery,newsizedata)
            return {"Newupdatedsize": newsize,"Unit": "MB","status":"success"}
      else:
        return {"status":"exceeds"}
